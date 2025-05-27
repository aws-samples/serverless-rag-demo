import boto3
from os import getenv
import xmltodict
import json
import datetime
import secrets
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime, timedelta
import os
import base64
from botocore.exceptions import ClientError 
import base64
from PIL import Image
import io
import os
from strands import Agent, tool
from agent_executor_utils import bedrockModel
import logging

date = datetime.now()
next_date = datetime.now() + timedelta(days=30)
year = date.year
month = date.month
month_label = date.strftime("%B")
next_month_label = next_date.strftime("%B")
day = date.day
bedrock_client = boto3.client('bedrock-runtime')
s3_client = boto3.client('s3')
s3_bucket_name = getenv("S3_BUCKET_NAME", "S3_BUCKET_NAME_MISSING")
cwd = os.getcwd() 
THEMES = ['strands/ppt_themes/ion_theme.pptx', 'strands/ppt_themes/circuit_theme.pptx']
cwd = getenv("LAMBDA_TASK_ROOT", 'var/')


PPT_GENERATOR_SYSTEM_PROMPT = """
You are a PPT Generator Agent. Your task is to generate a presentation on a given topic.
You have access to the following tools:
- generate_ppt: To generate a presentation on a given topic. 
                This tool will generate a presentation and upload it on s3.
                It will provide you the S3 key where the presentation is uploaded in <location>...</location> tags.

Key Responsibilities:
- Generate a presentation on a given topic
- Use the tools to generate a presentation

General Instructions:
- Use only the following slide formats: 
    a. Title page
    b. Slide with bullet points
    c. Slide with text,
    d. Slide with 4 takeaways

Slide specific instructions:
- For the Title page slide:
    - The first slide must be a Title page
    - The title should be a single line
    - The subtitle should be a single line
    - The text should be a single line
    - The speaker notes should be a single line
    - The slide format should be Title page

- For the intermediate slides:
    - The slide number should be incremented for each new slide
    - The title should be a single line
    - The subtitle should be a single line
    - The text should be a single line
    - The speaker notes should be a single line
    - The slide format should be one of the following:
        - Slide with bullet points
        - Slide with text
        - Slide with image and text
        - Slide with image only

- For the final Key Takeaways slide:    
    - The slide number should be incremented for each new slide
    - The title should be a single line
    - The subtitle should be a single line
    - The text should be a single line
    - The speaker notes should be a single line
    - The slide format should be Slide with 4 takeaways

- xml structure:
    - <presentation>
        - <slides>
            - <slide>
            - <slide_number>
            - <title>
            - <subtitle>
            - <text>
            - <speaker_notes>
            - <slideFormat>
            - </slide>
        - </slides>
    - </presentation>

Remember to:
- Use a unique and relevant title, subtitle, text, and speaker notes for each slide
- Vary the slide formats to make the presentation engaging
- Follow the specified xml structure

You should use the tools to generate a presentation.
You should also follow the slide specific instructions.
You should also follow the general instructions.
"""

LOG = logging.getLogger("ppt_generator_agent")
LOG.setLevel(logging.INFO)

@tool
def ppt_generator_agent(topic, additional_data_points, number_of_slides=3):
    agent = Agent(system_prompt=PPT_GENERATOR_SYSTEM_PROMPT, model=bedrockModel, tools=[generate_ppt])
    content = f"""Topic: {topic}
                Additional Data Points: {additional_data_points}
                Number of Slides: {number_of_slides}"""
        
    agent_response = agent(content)
    return agent_response


@tool
def generate_ppt(ppt_content):
    LOG.info(f"""Presentation Content: {ppt_content} """)
    LOG.debug(f"PPT prompt {PPT_GENERATOR_SYSTEM_PROMPT}")
    
    presentation = ""
    if '<presentation>' in ppt_content and '</presentation>' in ppt_content:
        presentation = ppt_content.split('</presentation')[0]
        presentation = presentation.split('<presentation>')[1]

    try:
        return ppt(presentation)
    except Exception as e:
        print(f"PPT generation error: {e}")
    return "PPT generation failed - 0"


def ppt(slides_xml_str: str):
    print(f"ppt_xml= {slides_xml_str}")
    slides_json = {}
    # convert xml to json
    slides_json = xmltodict.parse(slides_xml_str)
    print(f"xml parsed as dict={slides_json}")
    if 'slides' in slides_json:
        slides_json = slides_json['slides']
        if 'slide' in slides_json:
            slides_json = slides_json['slide']
    print(f"final_slides_json={slides_json}")
    theme = secrets.choice(THEMES)
    print(f"Theme selected: {theme}")
    ppt = Presentation(theme)
    # On local machine, use the following line to create a new presentation
    # ppt = Presentation()

    for slide_json in slides_json:
        if slide_json["slideFormat"] == "Title page":
            title_slide_layout = ppt.slide_layouts[0]
            slide = ppt.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = slide_json["title"]
            subtitle.text = slide_json["subtitle"]
        elif slide_json["slideFormat"] == "Slide with bullet points" or slide_json["slideFormat"]=="Slide with text":
            bullet_slide_layout = ppt.slide_layouts[1]
            slide = ppt.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = slide_json["title"]
            tf = body_shape.text_frame
            slide_json["text"] = str(slide_json["text"]).replace("*** Bullet points ***", '')
            tf.text = slide_json["text"]
            
        elif slide_json["slideFormat"] == "Slide with 4 takeaways":
            takeways_slide_layout = ppt.slide_layouts[1]
            slide = ppt.slides.add_slide(takeways_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            slide.placeholders[1].text = slide_json['subtitle']
            # For adjusting the  Margins in inches  
            left = width = height = Inches(2)
            top = Inches(4)
            # creating textBox 
            txBox = slide.shapes.add_textbox(left, top, width, height) 
            tf = txBox.text_frame
            p = tf.add_paragraph() 
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.italic = True
            p.text = slide_json["text"]
            
        elif slide_json["slideFormat"] == "Slide with image and text":
            # not done yet
            pass

    # img_path = "placeholder.png"
    print("PPT generated")
    now = datetime.now()
    date_time = now.strftime("%Y-%m-%d-%H-%M-%S")
    file_extension = 'pptx'
    file_name = f"/tmp/sample_{date_time}.{file_extension}"
    ppt.save(file_name)
    print("PPT Saved")
    print("PPT Upload in progress")
    is_upload, upload_location = upload_file_to_s3(file_name, file_extension)
    if(is_upload):
        return f"Presentation created successfully at location <location>{upload_location}</location>"
    else:
        return "PPT creation failed"

# Generate Images
def generate_bedrock_images(context, image_placeholder):
    height = 576
    width = 384

    system_prompt = f""" Summarize the following content in comma separated abstract concepts, maximum 20 words. 
            The text will be used to generate a representative image with Stable Diffusion. Remove preamble when anwering.
            """ + context
    model_id = 'amazon.titan-image-generator-v1'
    
    body = json.dumps({"taskType": "TEXT_IMAGE","textToImageParams": {
            "text": system_prompt
        },
        "imageGenerationConfig": {"numberOfImages": 1,
            "height": height, 
            "width": width,
            "cfgScale": 8.0, "seed": 1
        }
    })
    try:
        image_bytes = generate_image(model_id=model_id, body=body)
        image = Image.open(io.BytesIO(image_bytes))
        image.save(cwd+"/tmp/test_image.jpg")
        image_placeholder.insert_picture(cwd+'/tmp/test_image.jpg')
    except ClientError as err:
        message = err.response["Error"]["Message"]
        print("A client error occurred:", message)
        print("A client error occured: " + format(message))
    except Exception as err:
        print(err.message)
    else:
        print(f"Finished generating image with Amazon Titan Image Generator G1 model {model_id}.")

def generate_image(model_id, body):
    print("Generating image with Amazon Titan Image Generator G1 model", model_id)
    accept = "application/json"
    content_type = "application/json"

    response = bedrock_client.invoke_model(
        body=body, modelId=model_id, accept=accept, contentType=content_type
    )
    response_body = json.loads(response.get("body").read())

    base64_image = response_body.get("images")[0]
    base64_bytes = base64_image.encode('ascii')
    image_bytes = base64.b64decode(base64_bytes)

    finish_reason = response_body.get("error")

    if finish_reason is not None:
        print(f"Error Image generation error. Error is {finish_reason}")
        return None
    
    print("Successfully generated image with Amazon Titan Image Generator G1 model", model_id)
    return image_bytes

def upload_file_to_s3(file_name, file_extension):
    try:
        now = datetime.now()
        date_time = now.strftime("%Y-%m-%d-%H-%M-%S")
        s3_key = f"{file_extension}/sample_{date_time}.{file_extension}"
        s3_client.upload_file(file_name, s3_bucket_name, s3_key)
        s3_presigned = generate_presigned_url(s3_key)
        if s3_presigned is not None:
            return True, s3_presigned
        return True, s3_key
    except Exception as e:
        print(f"Error uploading to S3: {e}")
        return False, f'Error {e}'

# Generate a presigned get url for s3 file
def generate_presigned_url(s3_key):
    try:
        presigned_url = s3_client.generate_presigned_url(
            ClientMethod='get_object',
            Params={
                'Bucket': s3_bucket_name,
                'Key': s3_key
            },
            ExpiresIn=3600  # URL expires in 1 hour
        )
        return presigned_url
    except Exception as e:
        print(f"Error generating presigned URL: {e}")
        return None


# if __name__ == "__main__":
#     print(generate_ppt_agent("The Wonderful World of Cats", "Beautifull cats", 3))

# def test_ppt():
#     xml = """
# <slides>
# <slide>
# <slide_number>1</slide_number>
# <title>The Wonderful World of Cats</title>
# <subtitle>Explore the fascinating feline friends</subtitle>
# <text>This presentation will take a deep dive into the captivating lives of cats, covering their unique characteristics, behaviors, and the special bond they share with humans.</text>
# <speaker_notes>Welcome everyone! Today, we'll embark on a journey to discover the amazing world of our feline companions. Get ready to be enchanted by the fascinating facts and endearing traits of these adorable creatures.</speaker_notes>
# <slideFormat>Title page</slideFormat>
# </slide>
# <slide>
# <slide_number>2</slide_number>
# <title>Purr-fect Personalities</title>
# <subtitle>Exploring the diverse traits of cats</subtitle>
# <text>
# - Cats are known for their independence and self-reliance, often aloof yet affectionate
# - They exhibit a wide range of personalities, from playful and curious to calm and reserved
# - Cats are master communicators, using a variety of vocalizations and body language to express their needs and emotions
# - Each cat has its own unique quirks and idiosyncrasies that make them endearing companions
# </text>
# <speaker_notes>Cats are complex and multifaceted creatures, with a wide range of personalities and behaviors that make them truly captivating. Let's dive into some of the key traits that make cats such fascinating pets.</speaker_notes>
# <slideFormat>Slide with bullet points</slideFormat>
# </slide>
# <slide>
# <slide_number>3</slide_number>
# <title>The Feline Senses</title>
# <subtitle>Exploring the extraordinary abilities of cats</subtitle>
# <text>Cats possess a remarkable set of sensory capabilities that allow them to thrive in their environments:
# Sight: Cats have exceptional night vision and can see a wider range of colors than humans.
# Hearing: Feline ears can detect a wider range of frequencies, enabling them to hear sounds that humans cannot.
# Smell: A cat's sense of smell is up to 14 times more sensitive than a human's, allowing them to gather valuable information about their surroundings.
# These heightened senses contribute to the cat's impressive hunting prowess and overall adaptability.</text>
# <speaker_notes>Cats have evolved with an impressive array of sensory abilities that allow them to navigate their world with incredible precision and awareness. Let's explore some of the key ways in which their senses set them apart.</speaker_notes>
# <slideFormat>Slide with text</slideFormat>
# </slide>
# <slide>
# <slide_number>4</slide_number>
# <title>The Feline Agility</title>
# <subtitle>Marveling at the athletic prowess of cats</subtitle>
# <image>https://example.com/cat-jumping.jpg</image>
# <speaker_notes>Cats are renowned for their incredible agility and athleticism. From effortlessly leaping to great heights to executing precise landings, their physical capabilities are truly awe-inspiring.</speaker_notes>
# <slideFormat>Slide with image only</slideFormat>
# </slide>
# <slide>
# <slide_number>5</slide_number>
# <title>Important Takeaways</title>
# <subtitle>Key messages to remember</subtitle>
# <text>
# - Cats exhibit a wide range of unique personalities and traits
# - Feline senses are extraordinarily keen, enabling them to thrive in their environments
# - Cats possess remarkable agility and athleticism, allowing them to perform amazing feats
# - The special bond between cats and humans is a cherished and rewarding relationship
# </text>
# <speaker_notes>As we conclude our exploration of the wonderful world of cats, remember these key takeaways about these fascinating feline friends. Their captivating personalities, heightened senses, and incredible physical abilities make them truly remarkable creatures.</speaker_notes>
# <slideFormat>Slide with 4 takeaways</slideFormat>
# </slides>
# """
#     ppt(xml)

# test_ppt()