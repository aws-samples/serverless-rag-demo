import boto3
from os import getenv
import xmltodict
import json
import datetime
import secrets
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime, timedelta
import os
import base64
from botocore.exceptions import ClientError 
import base64
from PIL import Image
import io
import os
from agent_executor_utils import upload_file_to_s3

#Begin: Needed by Master orchestrator
ppt_agent_name = "PPT Generator Agent"

# When to use this agent
ppt_agent_uses = f""" 
If the user seeks to generate a presentation of a particular topic of interest, then use the {ppt_agent_name} to answer the question
Use the {ppt_agent_name} if:
   1. If data points are not provided, but the topic is known to you, you can come up with data points for the presentation
   2. If the data points are not provided and the topic is unknown to you, then ask the User for more data points
   3. If the user doesnt have enough information then you could check with other Agents and seek more information on a topic
"""

# Agent use examples
ppt_agent_use_examples = f"""
Create a presentation on AWS -> Use the {ppt_agent_name} to build the PPT
Generate a presentation on the above topic -> Use the {ppt_agent_name} to build the PPT
Build a presentation on the $TOPIC topic -> Use the {ppt_agent_name} to build the PPT
generate a ppt on the $TOPIC topic -> Use the {ppt_agent_name} to build the PPT
"""

# Agent success criteria
ppt_agent_stop_condition = f"""
  This {ppt_agent_name} agent successfully generates the PPT and uploads to S3 and shares the S3 key in the response
"""
# Agent Tool information
ppt_specs = f"""\
<agent_name>{ppt_agent_name}</agent_name>
<tool_set>
 <instructions>
   1. You are tasked to generate a powerpoint presentation based on the provided data points.
   2. Presentation Agent Rules:
        a. The Presentation Agent will create a presentation with the available context without requesting for more data points from the user
        b. The Presentation Agent will not ask any questions to the user
        c. The Presentation Agent will not write any explanations
           
   </instructions>

   <tool_description>
   <tool_usage>This tool is used to generate a PPT</tool_usage>
   <tool_name>generate_ppt</tool_name>
   <parameters>
   <parameter>
       <name>topic</name>
       <type>string</type>
       <description>Presentation TOPIC</description>
   </parameter>
   <parameter>
       <name>additional_data_points</name>
       <type>string</type>
       <description>Raw data points in the chat history that may be useful to create the ppt</description>
   </parameter>
   <parameter>
       <name>number_of_slides</name>
       <type>integer</type>
       <description>Total number of slides</description>
   </parameter>
   </parameters>
   </tool_description>   
<tool_set>
"""
#end: Needed by Master orchestrator

date = datetime.now()
next_date = datetime.now() + timedelta(days=30)
year = date.year
month = date.month
month_label = date.strftime("%B")
next_month_label = next_date.strftime("%B")
day = date.day


bedrock_client = boto3.client('bedrock-runtime')
model_id = getenv("PPT_MODEL", "anthropic.claude-3-sonnet-20240229-v1:0")
s3_bucket_name = getenv("S3_BUCKET_NAME", "S3_BUCKET_NAME_MISSING")
cwd = os.getcwd() 
THEMES = ['agents/artifacts/ion_theme.pptx', 'agents/artifacts/circuit_theme.pptx']
cwd = getenv("LAMBDA_TASK_ROOT", 'var/')

def generate_ppt(topic, additional_data_points="", number_of_slides=3):
    print(
        f"""\ Presentation Topic: {topic}, Number of Slides: {number_of_slides} """
    )

    system_prompt = f""" Your task is to propose a {number_of_slides} slide presentation on the topic: {topic}.
                    
                    The presentation should follow these requirements:
                    - Use only the following slide formats: 
                        a. Title page
                        b. Slide with bullet points
                        c. Slide with text,
                        d. Slide with 4 takeaways
                    - The first slide must be a Title page
                    - The last slide must be a Slide with 4 takeaways

                    Wrap the presentation XML within <presentation></presentation> tags.
                    Do not include any other text or tags in the response.
                    
            Your response should be in XML format with the following structure:
            <slides>
                <slide>
                   <slide_number>$SLIDE_NUMBER</slide_number>
                   <title>$TITLE</title>
                   <subtitle>$SUBTITLE</subtitle>
                   <text>$TEXT</text>
                   <speaker_notes>$SPEAKER_NOTES</speaker_notes>
                   <slideFormat>$SLIDE_FORMAT</slideFormat> 
                </slide>
            </slides>

        """          

    system_prompt = system_prompt +   f"""
      Below are some of the slide specific instructions:
        For the Title page slide:
        <slide_number>1</slide_number>
        <title>Come up with an engaging title for the presentation</title>
        <subtitle>Add a subtitle that captures the essence of the topic</subtitle>
        <text>Provide a brief overview of what the presentation will cover</text>
        <speaker_notes>Introduce yourself and give context for the presentation topic</speaker_notes>
        <slideFormat>Title page</slideFormat>
        
        For the intermediate slides:
        <slide_number>Increment this number for each new slide</slide_number>
        <title>Create a title summarizing the main point of this slide</title>
        <subtitle>Add a subtitle to complement the title</subtitle>
        <text>
        If using a Slide with bullet points format:
        *** Include 3-5 bullet points covering key information for this slide ***
        Else:
        Write 2-3 concise paragraphs with supporting details for the slide topic
        </text>
        <speaker_notes>Add relevant notes to help explain or expand on the slide content</speaker_notes>
        <slideFormat>
        Choose one of the following formats based on the content:
        - Slide with bullet points
        - Slide with image and text
        - Slide with image only
        </slideFormat>
        
        For the final Key Takeaways slide:
        <slide_number>{number_of_slides}</slide_number>
        <title>Important Takeaways</title>
        <subtitle>Key messages to remember</subtitle>
        <text>
        *** Summarize the 4 most crucial points covered in the presentation ***
        </text>
        <speaker_notes>Remind the audience of the key information you want them to walk away
        with</speaker_notes>
        <slideFormat>Slide with 4 takeaways</slideFormat>
        
        Remember to:
        - Use a unique and relevant title, subtitle, text, and speaker notes for each slide
        - Vary the slide formats to make the presentation engaging
        - Follow the specified xml structure
        """
    
    print(f"PPT prompt {system_prompt}")
    user_context = "Generate a presentation json"
    if additional_data_points != "":
        user_context = f"""The following context was provided by the user:
                        <context>
                        {additional_data_points} </context> 
                        <instructions>
                        When creating the presentation json you should consider the context available in <context></context> tags
                        You should generate a well formatted accurate presentation xml
                        <instructions>
                        """
        
    query_list = [{"role": "user", "content":  user_context }]
    
    prompt_template = {
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 10000,
        "system": system_prompt,
        "messages": query_list
    }
    
    # print(f'Prompt Template = {prompt_template}')
    
    response = bedrock_client.invoke_model(
        body=json.dumps(prompt_template),
        modelId=model_id,
        accept="application/json",
        contentType="application/json",
    )
    llm_output = json.loads(response["body"].read())
    presentation = ""
    if "content" in llm_output:
        presentation = llm_output["content"][0]["text"]
        if '<presentation>' in presentation and '</presentation>' in presentation:
            presentation = presentation.split('</presentation')[0]
            presentation = presentation.split('<presentation>')[1]

    try:
        return ppt(presentation)
    except Exception as e:
        print(f"PPT generation error: {e}")
    return "PPT generation failed - 0"

def ppt(slides_xml_str: str):
    print(f"ppt_xml= {slides_xml_str}")
    slides_json = {}
    # convert xml to json
    slides_json = xmltodict.parse(slides_xml_str)
    print(f"xml parsed as dict={slides_json}")
    if 'slides' in slides_json:
        slides_json = slides_json['slides']
        if 'slide' in slides_json:
            slides_json = slides_json['slide']
    print(f"final_slides_json={slides_json}")
    theme = secrets.choice(THEMES)
    print(f"Theme selected: {theme}")
    ppt = Presentation(theme)

    for slide_json in slides_json:
        if slide_json["slideFormat"] == "Title page":
            title_slide_layout = ppt.slide_layouts[0]
            slide = ppt.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = slide_json["title"]
            subtitle.text = slide_json["subtitle"]
        elif slide_json["slideFormat"] == "Slide with bullet points" or slide_json["slideFormat"]=="Slide with text":
            bullet_slide_layout = ppt.slide_layouts[1]
            slide = ppt.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            title_shape.text = slide_json["title"]
            tf = body_shape.text_frame
            slide_json["text"] = str(slide_json["text"]).replace("*** Bullet points ***", '')
            tf.text = slide_json["text"]
            
        elif slide_json["slideFormat"] == "Slide with 4 takeaways":
            takeways_slide_layout = ppt.slide_layouts[1]
            slide = ppt.slides.add_slide(takeways_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            slide.placeholders[1].text = slide_json['subtitle']
            # For adjusting the  Margins in inches  
            left = width = height = Inches(2)
            top = Inches(4)
            # creating textBox 
            txBox = slide.shapes.add_textbox(left, top, width, height) 
            tf = txBox.text_frame
            p = tf.add_paragraph() 
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.italic = True
            p.text = slide_json["text"]
            
        elif slide_json["slideFormat"] == "Slide with image and text":
            # not done yet
            pass

    # img_path = "placeholder.png"
    print("PPT generated")
    now = datetime.now()
    date_time = now.strftime("%Y-%m-%d-%H-%M-%S")
    file_extension = 'pptx'
    file_name = f"/tmp/sample_{date_time}.{file_extension}"
    ppt.save(file_name)
    print("PPT Saved")
    print("PPT Upload in progress")
    is_upload, upload_location = upload_file_to_s3(file_name, file_extension)
    if(is_upload):
        return f"Presentation created successfully at location <location>{upload_location}</location>"
    else:
        return "PPT creation failed"

# Generate Images
def generate_bedrock_images(context, image_placeholder):
    height = 576
    width = 384

    system_prompt = f""" Summarize the following content in comma separated abstract concepts, maximum 20 words. 
            The text will be used to generate a representative image with Stable Diffusion. Remove preamble when anwering.
            """ + context
    model_id = 'amazon.titan-image-generator-v1'
    
    body = json.dumps({"taskType": "TEXT_IMAGE","textToImageParams": {
            "text": system_prompt
        },
        "imageGenerationConfig": {"numberOfImages": 1,
            "height": height, 
            "width": width,
            "cfgScale": 8.0, "seed": 1
        }
    })
    try:
        image_bytes = generate_image(model_id=model_id, body=body)
        image = Image.open(io.BytesIO(image_bytes))
        image.save(cwd+"/tmp/test_image.jpg")
        image_placeholder.insert_picture(cwd+'/tmp/test_image.jpg')
    except ClientError as err:
        message = err.response["Error"]["Message"]
        print("A client error occurred:", message)
        print("A client error occured: " + format(message))
    except Exception as err:
        print(err.message)
    else:
        print(f"Finished generating image with Amazon Titan Image Generator G1 model {model_id}.")

def generate_image(model_id, body):
    print("Generating image with Amazon Titan Image Generator G1 model", model_id)

    accept = "application/json"
    content_type = "application/json"

    response = bedrock_client.invoke_model(
        body=body, modelId=model_id, accept=accept, contentType=content_type
    )
    response_body = json.loads(response.get("body").read())

    base64_image = response_body.get("images")[0]
    base64_bytes = base64_image.encode('ascii')
    image_bytes = base64.b64decode(base64_bytes)

    finish_reason = response_body.get("error")

    if finish_reason is not None:
        print(f"Error Image generation error. Error is {finish_reason}")
        return None
    
    print("Successfully generated image with Amazon Titan Image Generator G1 model", model_id)
    return image_bytes


# def upload_to_s3(file_name):
#     try:
#         s3_client = boto3.client("s3")
#         now = datetime.now()
#         date_time = now.strftime("%Y-%m-%d-%H-%M-%S")
#         s3_key = f"pptx/sample_{date_time}.pptx"
#         bucket_name = s3_bucket_name
#         s3_client.upload_file(file_name, bucket_name, s3_key)
#         return True, s3_key
#     except Exception as e:
        # print(f"Error uploading to S3: {e}")
        # return False, ''

def test_ppt():
    xml = """
<slides>
<slide>
<slide_number>1</slide_number>
<title>The Wonderful World of Cats</title>
<subtitle>Explore the fascinating feline friends</subtitle>
<text>This presentation will take a deep dive into the captivating lives of cats, covering their unique characteristics, behaviors, and the special bond they share with humans.</text>
<speaker_notes>Welcome everyone! Today, we'll embark on a journey to discover the amazing world of our feline companions. Get ready to be enchanted by the fascinating facts and endearing traits of these adorable creatures.</speaker_notes>
<slideFormat>Title page</slideFormat>
</slide>
<slide>
<slide_number>2</slide_number>
<title>Purr-fect Personalities</title>
<subtitle>Exploring the diverse traits of cats</subtitle>
<text>
- Cats are known for their independence and self-reliance, often aloof yet affectionate
- They exhibit a wide range of personalities, from playful and curious to calm and reserved
- Cats are master communicators, using a variety of vocalizations and body language to express their needs and emotions
- Each cat has its own unique quirks and idiosyncrasies that make them endearing companions
</text>
<speaker_notes>Cats are complex and multifaceted creatures, with a wide range of personalities and behaviors that make them truly captivating. Let's dive into some of the key traits that make cats such fascinating pets.</speaker_notes>
<slideFormat>Slide with bullet points</slideFormat>
</slide>
<slide>
<slide_number>3</slide_number>
<title>The Feline Senses</title>
<subtitle>Exploring the extraordinary abilities of cats</subtitle>
<text>Cats possess a remarkable set of sensory capabilities that allow them to thrive in their environments:
Sight: Cats have exceptional night vision and can see a wider range of colors than humans.
Hearing: Feline ears can detect a wider range of frequencies, enabling them to hear sounds that humans cannot.
Smell: A cat's sense of smell is up to 14 times more sensitive than a human's, allowing them to gather valuable information about their surroundings.
These heightened senses contribute to the cat's impressive hunting prowess and overall adaptability.</text>
<speaker_notes>Cats have evolved with an impressive array of sensory abilities that allow them to navigate their world with incredible precision and awareness. Let's explore some of the key ways in which their senses set them apart.</speaker_notes>
<slideFormat>Slide with text</slideFormat>
</slide>
<slide>
<slide_number>4</slide_number>
<title>The Feline Agility</title>
<subtitle>Marveling at the athletic prowess of cats</subtitle>
<image>https://example.com/cat-jumping.jpg</image>
<speaker_notes>Cats are renowned for their incredible agility and athleticism. From effortlessly leaping to great heights to executing precise landings, their physical capabilities are truly awe-inspiring.</speaker_notes>
<slideFormat>Slide with image only</slideFormat>
</slide>
<slide>
<slide_number>5</slide_number>
<title>Important Takeaways</title>
<subtitle>Key messages to remember</subtitle>
<text>
- Cats exhibit a wide range of unique personalities and traits
- Feline senses are extraordinarily keen, enabling them to thrive in their environments
- Cats possess remarkable agility and athleticism, allowing them to perform amazing feats
- The special bond between cats and humans is a cherished and rewarding relationship
</text>
<speaker_notes>As we conclude our exploration of the wonderful world of cats, remember these key takeaways about these fascinating feline friends. Their captivating personalities, heightened senses, and incredible physical abilities make them truly remarkable creatures.</speaker_notes>
<slideFormat>Slide with 4 takeaways</slideFormat>
</slides>
"""
    ppt(xml)

# test_ppt()