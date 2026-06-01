from strands import Agent
from strands.models import BedrockModel
from pptx import Presentation
from pptx.util import Pt
import json
import os
import io
import sys
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from utils import upload_and_get_url

MODEL_ID = os.getenv("MODEL_ID", "global.anthropic.claude-sonnet-4-6-v1:0")
REGION = os.getenv("REGION", "us-east-1")

PPT_PROMPT = """You are a presentation content planner. Given a topic, create slide content as JSON:
{"slides": [{"title": "...", "bullets": ["...", "..."]}]}
Generate 5-8 slides with clear, concise bullet points. Return ONLY valid JSON."""


def generate_ppt(query: str) -> str:
    model = BedrockModel(model_id=MODEL_ID, region_name=REGION)
    agent = Agent(system_prompt=PPT_PROMPT, model=model)
    response = str(agent(query))

    try:
        slides_data = json.loads(response)
    except json.JSONDecodeError:
        return "Failed to generate presentation structure."

    prs = Presentation()
    for slide_data in slides_data.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data.get("title", "")
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for bullet in slide_data.get("bullets", []):
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    url = upload_and_get_url(
        buffer.read(), "generated-ppt", "pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation")
    return f"Presentation generated. <location>{url}</location>"
